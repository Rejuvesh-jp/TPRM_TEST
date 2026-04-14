"""
OCR Service
===========
Extracts text from documents, with OCR fallback for scanned PDFs.
Supports: PDF, DOCX, TXT, CSV, XLSX, JSON
"""

import io
import logging
from pathlib import Path
import time
import platform
import tempfile
import shutil
import subprocess

logger = logging.getLogger("tprm.ocr_service")

# ── OCR text cache (keyed by SHA-256 of file bytes) ────────────────
# Prevents re-running expensive OCR/extraction on the same file bytes
# across repeated pipeline runs, guaranteeing identical extracted text.
_ocr_cache: dict[str, str] = {}


def _file_hash(file_path: Path) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def _get_cached_text(file_path: Path) -> str | None:
    """Return cached extracted text for a file if its content hash is known."""
    try:
        key = _file_hash(file_path)
        return _ocr_cache.get(key)
    except Exception:
        return None


def _put_cached_text(file_path: Path, text: str) -> None:
    """Store extracted text keyed by file content hash."""
    try:
        key = _file_hash(file_path)
        _ocr_cache[key] = text
    except Exception:
        pass


def extract_text_from_pdf(file_path: str | Path) -> str:
    """Extract text from a PDF. Falls back to OCR if text layer is empty."""
    from pypdf import PdfReader

    reader = PdfReader(str(file_path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i + 1}]\n{text}")

    if pages:
        return "\n\n".join(pages)

    # Fallback: attempt OCR via pytesseract + pdf2image
    logger.info("No text layer in %s — attempting OCR", file_path)
    return _ocr_pdf(file_path)


def _ocr_pdf(file_path: str | Path) -> str:
    """
    Cross-platform scanned PDF OCR
    Windows  -> Windows OCR
    Linux/Mac -> Tesseract
    """
    import fitz
    import time

    system = platform.system()

    try:
        doc = fitz.open(str(file_path))
        mat = fitz.Matrix(200 / 72, 200 / 72)

        with tempfile.TemporaryDirectory() as tmp_dir:
            img_paths = []
            t_render = time.perf_counter()
            for i in range(len(doc)):
                pix = doc[i].get_pixmap(matrix=mat)
                img_path = Path(tmp_dir) / f"page_{i}.png"
                pix.save(str(img_path))
                img_paths.append(img_path)

            render_ms = (time.perf_counter() - t_render) * 1000
            logger.info(
                "[PERF] OCR render %d page(s): %.0f ms",
                len(img_paths),
                render_ms
            )

            doc.close()

            extracted_pages = []

            for img_path in img_paths:
                if system == "Windows":
                    text = _ocr_windows_image(img_path)
                else:
                    text = _ocr_tesseract_image(img_path)

                if text:
                    extracted_pages.append(text)

            return "\n".join(extracted_pages)

    except Exception as exc:
        logger.warning("PDF OCR failed for %s: %s", file_path, exc)
        return ""

def _build_ocr_ps_script(image_paths: list[str]) -> str:
    """Build a PowerShell script that uses Windows.Media.Ocr to OCR images."""
    # Escape backslashes for PowerShell string literals
    paths_array = ", ".join(f'"{p}"' for p in image_paths)

    return f'''
Add-Type -AssemblyName "System.Runtime.WindowsRuntime"
[Windows.Media.Ocr.OcrEngine, Windows.Foundation.UniversalApiContract, ContentType = WindowsRuntime] | Out-Null
[Windows.Graphics.Imaging.BitmapDecoder, Windows.Foundation.UniversalApiContract, ContentType = WindowsRuntime] | Out-Null
[Windows.Storage.StorageFile, Windows.Storage, ContentType = WindowsRuntime] | Out-Null

$Async = ([System.WindowsRuntimeSystemExtensions].GetMethods() | Where-Object {{
    $_.Name -eq 'AsTask' -and $_.GetParameters().Count -eq 1 -and $_.GetParameters()[0].ParameterType.Name -eq 'IAsyncOperation`1'
}})[0]

Function Await($AsyncOp, [Type] $ResultType) {{
    $Async.MakeGenericMethod($ResultType).Invoke($null, @($AsyncOp)).GetAwaiter().GetResult()
}}

$ocrEngine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromUserProfileLanguages()
if ($null -eq $ocrEngine) {{ exit 1 }}

$imagePaths = @({paths_array})

foreach ($imgPath in $imagePaths) {{
    try {{
        $file = Await ([Windows.Storage.StorageFile]::GetFileFromPathAsync($imgPath)) ([Windows.Storage.StorageFile])
        $stream = Await ($file.OpenAsync([Windows.Storage.FileAccessMode]::Read)) ([Windows.Storage.Streams.IRandomAccessStream])
        $decoder = Await ([Windows.Graphics.Imaging.BitmapDecoder]::CreateAsync($stream)) ([Windows.Graphics.Imaging.BitmapDecoder])
        $bitmap = Await ($decoder.GetSoftwareBitmapAsync()) ([Windows.Graphics.Imaging.SoftwareBitmap])
        $ocrResult = Await ($ocrEngine.RecognizeAsync($bitmap)) ([Windows.Media.Ocr.OcrResult])

        foreach ($line in $ocrResult.Lines) {{
            Write-Output $line.Text
        }}

        $stream.Dispose()
        $bitmap.Dispose()
    }} catch {{
        # Skip failed page
    }}
}}
'''


def extract_text_from_docx(file_path: str | Path) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text_from_txt(file_path: str | Path) -> str:
    """Extract text from a plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_text_from_csv(file_path: str | Path) -> str:
    """Extract text from a CSV file."""
    import csv

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = [" | ".join(row) for row in reader if any(row)]
    return "\n".join(rows)


def extract_text_from_xlsx(file_path: str | Path) -> str:
    """Extract text from an Excel file."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def extract_text_from_pptx(file_path: str | Path) -> str:
    """Extract text from a PowerPoint (.pptx) file."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_text_from_xls(file_path: str | Path) -> str:
    """Extract text from a legacy Excel (.xls) file."""
    import xlrd

    wb = xlrd.open_workbook(str(file_path))
    sheets = []
    for sheet in wb.sheets():
        rows = []
        for rx in range(sheet.nrows):
            cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def extract_text_from_json(file_path: str | Path) -> str:
    """Extract text from a JSON file."""
    import json as _json

    with open(file_path, "r", encoding="utf-8") as f:
        data = _json.load(f)
    return _json.dumps(data, indent=2)


def extract_text_from_image(file_path: str | Path) -> str:
    """
    Cross-platform OCR for images.
    Windows  -> Windows OCR API
    Linux/Mac -> Tesseract OCR
    """
    file_path = Path(file_path)

    system = platform.system()

    try:
        if system == "Windows":
            return _ocr_windows_image(file_path)

        elif system in ("Linux", "Darwin"):
            return _ocr_tesseract_image(file_path)

        else:
            logger.warning("Unsupported platform for OCR: %s", system)
            return ""

    except Exception as exc:
        logger.warning("Image OCR failed for %s: %s", file_path.name, exc)
        return ""


def _ocr_tesseract_image(file_path: Path) -> str:
    """OCR image using Tesseract (Linux/macOS)."""
    if not shutil.which("tesseract"):
        logger.warning("Tesseract is not installed. Skipping OCR for %s", file_path)
        return ""

    result = subprocess.run(
        ["tesseract", str(file_path), "stdout"],
        capture_output=True,
        text=True,
        timeout=120
    )

    if result.returncode == 0:
        text = result.stdout.strip()
        logger.info("OCR extracted %d chars from image %s", len(text), file_path.name)
        return text

    logger.warning("Tesseract OCR failed for %s: %s", file_path, result.stderr[:500])
    return ""


def _ocr_windows_image(file_path: Path) -> str:
    """OCR image using Windows OCR."""
    abs_path = str(file_path.resolve())
    ps_script = _build_ocr_ps_script([abs_path])

    with tempfile.TemporaryDirectory() as tmp_dir:
        ps_file = Path(tmp_dir) / "ocr_img.ps1"
        ps_file.write_text(ps_script, encoding="utf-8")

        result = subprocess.run(
            ["powershell", "-ExecutionPolicy", "Bypass", "-File", str(ps_file)],
            capture_output=True,
            text=True,
            timeout=120
        )

        return result.stdout.strip() if result.returncode == 0 else ""
# ── Dispatcher ─────────────────────────────────────────

EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".txt": extract_text_from_txt,
    ".csv": extract_text_from_csv,
    ".xlsx": extract_text_from_xlsx,
    ".xls": extract_text_from_xls,
    ".pptx": extract_text_from_pptx,
    ".json": extract_text_from_json,
    ".png": extract_text_from_image,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
}

SUPPORTED_EXTENSIONS = set(EXTRACTORS.keys())


import re as _re

# Mapping of Unicode characters commonly mangled during PDF/OCR extraction
_NORMALIZE_MAP = {
    "\u2018": "'",   # LEFT SINGLE QUOTATION MARK  → apostrophe
    "\u2019": "'",   # RIGHT SINGLE QUOTATION MARK → apostrophe
    "\u201A": "'",   # SINGLE LOW-9 QUOTATION MARK → apostrophe
    "\u201B": "'",   # SINGLE HIGH-REVERSED-9       → apostrophe
    "\u201C": '"',   # LEFT DOUBLE QUOTATION MARK   → quote
    "\u201D": '"',   # RIGHT DOUBLE QUOTATION MARK  → quote
    "\u201E": '"',   # DOUBLE LOW-9 QUOTATION MARK  → quote
    "\u2013": "-",   # EN DASH  → hyphen
    "\u2014": "-",   # EM DASH  → hyphen
    "\u2026": "...", # HORIZONTAL ELLIPSIS → three dots
    "\u00A0": " ",   # NO-BREAK SPACE → regular space
    "\uFEFF": "",    # BOM / ZERO WIDTH NO-BREAK SPACE → remove
}
_NORMALIZE_RE = _re.compile("|".join(_re.escape(k) for k in _NORMALIZE_MAP))


def _normalize_text(text: str) -> str:
    """Normalize smart quotes, fancy dashes, and other Unicode artifacts from PDF/OCR extraction."""
    if not text:
        return text
    # Strip NULL bytes and other control characters that PostgreSQL cannot store
    text = text.replace("\x00", "")
    text = _re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    # Replace known Unicode characters
    text = _NORMALIZE_RE.sub(lambda m: _NORMALIZE_MAP[m.group()], text)
    # Fix orphaned spaces before common possessive/contraction suffixes
    # e.g. "Provider s" → "Provider's", "don t" → "don't", "we re" → "we're"
    text = _re.sub(r"(?<=\w) (?=s\b(?:[^a-zA-Z]|$))", "'", text)
    text = _re.sub(r"\b(don|won|can|shouldn|wouldn|couldn|isn|aren|wasn|weren|hasn|haven|hadn|didn|doesn|ain) (t)\b",
                   r"\1'\2", text, flags=_re.IGNORECASE)
    text = _re.sub(r"\b(I|you|we|they|who) (m|re|ll|ve|d)\b",
                   r"\1'\2", text, flags=_re.IGNORECASE)
    text = _re.sub(r"\b(it|he|she|that|there|here|what|where|how) (s|ll|d)\b",
                   r"\1'\2", text, flags=_re.IGNORECASE)
    return text


def extract_text(file_path: str | Path) -> str:
    """Extract text from a file based on its extension.

    Results are cached by file content hash so identical files are never
    re-processed — this is the primary guard against OCR output drift.
    """
    file_path = Path(file_path)
    cached = _get_cached_text(file_path)
    if cached is not None:
        logger.debug("OCR cache HIT: %s", file_path.name)
        return cached

    ext = file_path.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")
    logger.info("Extracting text from %s (type: %s)", file_path, ext)
    raw_text = extractor(file_path)
    result = _normalize_text(raw_text)
    _put_cached_text(file_path, result)
    return result
